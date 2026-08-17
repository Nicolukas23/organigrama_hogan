from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

RED = RGBColor(0xE8, 0x11, 0x2D)
DARK = RGBColor(0x1A, 0x1A, 0x2E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x66, 0x66, 0x66)
GREEN = RGBColor(0x27, 0xAE, 0x60)
ORANGE = RGBColor(0xF3, 0x9C, 0x12)
BLUE = RGBColor(0x29, 0x80, 0xB9)
CARD = RGBColor(0x2A, 0x2A, 0x40)

def add_bg(slide, color=DARK):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def add_text(slide, l, t, w, h, text, sz=14, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(sz); p.font.bold = bold; p.font.color.rgb = color; p.font.name = 'Calibri'
    p.alignment = align
    return tb

def add_box(slide, l, t, w, h, fill=CARD, text='', sz=12, fc=WHITE, bold=False, align=PP_ALIGN.CENTER):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill; sh.line.fill.background()
    if text:
        tf = sh.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = text
        p.font.size = Pt(sz); p.font.color.rgb = fc; p.font.bold = bold; p.font.name = 'Calibri'
        p.alignment = align
    return sh

def add_bullets(slide, l, t, w, h, items, sz=13, color=WHITE):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item; p.font.size = Pt(sz); p.font.color.rgb = color; p.font.name = 'Calibri'
        p.space_after = Pt(6)

def section_divider(title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK)
    add_box(slide, 5, 3, 3.333, 0.06, RED)
    add_text(slide, 1, 2, 11.333, 1, title, 36, True, RED, PP_ALIGN.CENTER)
    add_text(slide, 1, 3.3, 11.333, 0.6, subtitle, 16, False, GRAY, PP_ALIGN.CENTER)
    return slide

# =============================================
# PORTADA
# =============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_box(slide, 0, 0, 13.333, 0.08, RED)
add_text(slide, 1, 1.5, 11, 1, 'ASSESSMENT CENTER', 42, True, RED, PP_ALIGN.CENTER)
add_text(slide, 1, 2.5, 11, 1, 'Analista Senior / Data Analytics', 28, False, WHITE, PP_ALIGN.CENTER)
add_text(slide, 1, 3.8, 11, 0.6, 'Análisis Exploratorio de Dataset Organizacional', 18, False, GRAY, PP_ALIGN.CENTER)
add_box(slide, 4, 5, 5.333, 0.06, RED)
add_text(slide, 1, 5.5, 11, 0.5, 'Bloque 1 — Caso Técnico (90 min)', 16, False, GRAY, PP_ALIGN.CENTER)

# =============================================
# BLOQUE 2: PRESENTACIÓN EJECUTIVA
# =============================================
section_divider('BLOQUE 2', 'Presentación Ejecutiva — 20 min — VP de Operaciones + Director de Talento')

# --- Slide: Estado de los datos ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_box(slide, 0, 0, 13.333, 0.08, RED)
add_text(slide, 0.5, 0.3, 12, 0.6, 'ESTADO DE LOS DATOS', 28, True, RED)
add_text(slide, 0.5, 0.9, 12, 0.5, 'Dataset: 15.000 registros × 10 columnas', 16, False, GRAY)

kpis = [('15.000', 'Registros'), ('10', 'Columnas'), ('98%', 'IDs válidos'), ('81', 'Salarios\n"Inconsistente"'), ('69', 'Salarios\nnegativos'), ('26', 'IDs\nduplicados')]
for i, (n, l) in enumerate(kpis):
    x = 0.5 + i * 2.1
    add_box(slide, x, 1.6, 1.9, 1.3, CARD)
    add_text(slide, x, 1.7, 1.9, 0.6, n, 28, True, RED, PP_ALIGN.CENTER)
    add_text(slide, x, 2.2, 1.9, 0.6, l, 11, False, GRAY, PP_ALIGN.CENTER)

issues = [
    'Columna "Segumiento" con error tipográfico',
    '150 registros (1%) sin Área definida',
    '300 registros (2%) sin ID de identificación',
    'Áreas duplicadas: "Compras" (1.342) vs "Comprass" (66)',
    '81 registros con Salario = "Inconsistente"',
    '69 salarios negativos',
    '1.373 outliers salariales (método IQR)',
    '"Corporativo" tiene solo 1 persona (Presidente) con $99,6M',
]
add_bullets(slide, 0.5, 3.2, 12, 4, ['• ' + i for i in issues], 13, WHITE)

# --- Slide: Hallazgos 1-2 ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_box(slide, 0, 0, 13.333, 0.08, RED)
add_text(slide, 0.5, 0.3, 12, 0.6, 'PRINCIPALES HALLAZGOS', 28, True, RED)

add_box(slide, 0.5, 1.1, 6, 0.5, RED, 'HALLAZGO 1: Distribución por Área', 14, WHITE, True)
areas = [('CAV','1.417','9,4%'),('Cartera','1.364','9,1%'),('Operación','1.348','9,0%'),('Compras','1.342','8,9%'),('IT','1.331','8,9%'),('Ventas','1.329','8,9%'),('Jurídica','1.321','8,8%'),('Marketing','1.293','8,6%'),('Servicio','1.282','8,5%'),('Tropas','1.261','8,4%'),('Ingeniería','1.261','8,4%')]
y = 1.7
for a, c, p in areas:
    add_text(slide, 0.5, y, 2, 0.25, a, 11, a=='CAV', ORANGE if a=='CAV' else WHITE)
    add_text(slide, 2.5, y, 1.2, 0.25, c, 11, False, WHITE)
    add_text(slide, 3.7, y, 1, 0.25, p, 11, False, GRAY)
    y += 0.25

add_box(slide, 6.8, 1.1, 6, 0.5, RED, 'HALLAZGO 2: Distribución por Cargo', 14, WHITE, True)
cargos = [('Analista','13.499','90,0%'),('Supervisor','450','3,0%'),('Director','300','2,0%'),('Jefe','300','2,0%'),('Gerente','300','2,0%'),('Director Ejecutivo','150','1,0%'),('Presidente','1','<0,1%')]
y = 1.7
for c, n, p in cargos:
    add_text(slide, 6.8, y, 2.5, 0.25, c, 11, c=='Presidente', RED if c=='Presidente' else WHITE)
    add_text(slide, 9.3, y, 1.2, 0.25, n, 11, False, WHITE)
    add_text(slide, 10.5, y, 1, 0.25, p, 11, False, GRAY)
    y += 0.25

add_box(slide, 6.8, 4.2, 6, 2.5, CARD)
add_text(slide, 7, 4.3, 5.6, 0.4, 'INSIGHT CLAVE', 14, True, ORANGE)
insights = ['• 90% son Analistas — pirámide muy plana','• Solo 1 Presidente en "Corporativo" ($99,6M)','• Relación Analista:Supervisor = 30:1','• CAV es el área más grande (1.417 personas)']
add_bullets(slide, 7, 4.7, 5.6, 1.8, insights, 12, WHITE)

# --- Slide: Hallazgos 3-4 ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_box(slide, 0, 0, 13.333, 0.08, RED)
add_text(slide, 0.5, 0.3, 12, 0.6, 'PRINCIPALES HALLAZGOS (cont.)', 28, True, RED)

add_box(slide, 0.5, 1.1, 4, 0.5, RED, 'HALLAZGO 3: Equilibrio de Género', 14, WHITE, True)
add_box(slide, 0.5, 1.7, 4, 2, CARD)
add_text(slide, 0.7, 1.8, 3.6, 0.5, 'Hombres: 7.502 (50,0%)', 16, True, BLUE, PP_ALIGN.CENTER)
add_text(slide, 0.7, 2.3, 3.6, 0.5, 'Mujeres: 7.498 (50,0%)', 16, True, RED, PP_ALIGN.CENTER)
add_text(slide, 0.7, 2.9, 3.6, 0.5, 'Paridad perfecta', 12, True, GREEN, PP_ALIGN.CENTER)

add_box(slide, 4.8, 1.1, 4, 0.5, RED, 'HALLAZGO 4: Área CAV', 14, WHITE, True)
add_box(slide, 4.8, 1.7, 4, 2, CARD)
cav = ['• Total: 1.417 personas (9,4%)','• 1.273 Analistas, 40 Supervisores','• 32 Gerentes, 29 Directores','• Salario promedio: $9.940.872']
add_bullets(slide, 5, 1.8, 3.6, 1.8, cav, 11, WHITE)

add_box(slide, 9.1, 1.1, 4, 0.5, RED, 'HALLAZGO 5: Análisis Salarial', 14, WHITE, True)
add_box(slide, 9.1, 1.7, 4, 3, CARD)
sal = ['• Promedio: $9.761.535','• Mediana: $6.394.501','• Mínimo: -$4.983.451 (negativo)','• Máximo: $99.620.396','• 81 "Inconsistente"','• 69 negativos','• 1.373 outliers (IQR)','⚠ Los salarios NO son confiables']
add_bullets(slide, 9.3, 1.8, 3.6, 2.8, sal, 11, WHITE)

# --- Slide: Riesgos ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_box(slide, 0, 0, 13.333, 0.08, RED)
add_text(slide, 0.5, 0.3, 12, 0.6, 'RIESGOS DE USAR ESTOS DATOS', 28, True, RED)

risks = [
    ('CRÍTICO', 'Salarios corruptos', '81 "Inconsistente" + 69 negativos + 1.373 outliers', RED),
    ('CRÍTICO', 'Datos faltantes', '300 sin ID + 150 sin Área = registros incompletos', RED),
    ('ALTO', 'Inconsistencia en áreas', '"Compras" vs "Comprass" + "undefined"', ORANGE),
    ('ALTO', 'IDs duplicados', '26 personas con mismo ID — distorsiona métricas', ORANGE),
    ('MEDIO', 'Columna "Seguimiento" vacía', 'Solo 107 de 15.000 (0,7%)', GREEN),
    ('MEDIO', 'Sin validación cruzada', 'Sexo = Sexo Cargado al 100%', GREEN),
]
y = 1.2
for tag, title, desc, color in risks:
    add_box(slide, 0.5, y, 12.3, 0.85, CARD)
    add_text(slide, 0.7, y+0.05, 1.5, 0.3, tag, 11, True, color)
    add_text(slide, 2.2, y+0.05, 3, 0.3, title, 14, True, WHITE)
    add_text(slide, 2.2, y+0.35, 10, 0.4, desc, 12, False, GRAY)
    y += 0.95

add_box(slide, 0.5, 7.0, 12.3, 0.5, RED)
add_text(slide, 0.7, 7.05, 12, 0.4, 'VEREDICTO: Sí se puede usar, PERO primero hay que limpiar los datos.', 14, True, WHITE, PP_ALIGN.CENTER)

# --- Slide: Recomendaciones ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_box(slide, 0, 0, 13.333, 0.08, RED)
add_text(slide, 0.5, 0.3, 12, 0.6, 'RECOMENDACIONES ACCIONABLES', 28, True, RED)

add_box(slide, 0.5, 1.1, 6, 0.5, RED, 'ACCIONES INMEDIATAS', 14, WHITE, True)
inmediatas = ['1. Limpiar salarios: "Inconsistente" → NULL, eliminar negativos','2. Resolver 300 IDs faltantes (JOIN con RRHH)','3. Unificar áreas: Compras/Comprass, reclasificar undefined','4. Eliminar 26 IDs duplicados','5. Corregir "Segumiento" → "Seguimiento"']
add_bullets(slide, 0.5, 1.7, 6, 3, inmediatas, 12, WHITE)

add_box(slide, 6.8, 1.1, 6, 0.5, RED, 'ACCIONES DE TRANSFORMACIÓN', 14, WHITE, True)
transform = ['6. Crear campo "Salario_Limpio" = válido o NULL','7. Agregar scoring de calidad: Completo/Parcial/Crítico','8. Validar contra fuente SAPHIRE (nómina real)','9. Regla: salario no puede ser negativo ni >3x promedio','10. Auditoría periódica de calidad']
add_bullets(slide, 6.8, 1.7, 6, 3, transform, 12, WHITE)

# =============================================
# BLOQUE 3: STORYTELLING (1 SLIDE)
# =============================================
section_divider('BLOQUE 3', 'Story Telling — 30 min — 1 slide para VP (15 min)')

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_box(slide, 0, 0, 13.333, 0.08, RED)
add_text(slide, 0.5, 0.2, 12, 0.5, 'RESUMEN EJECUTIVO PARA VP', 28, True, RED)
add_text(slide, 0.5, 0.7, 12, 0.4, 'Dataset de 15.000 empleados — ¿Puede usarlo para decisiones estratégicas?', 16, False, GRAY)

add_box(slide, 0.5, 1.3, 12.3, 1, CARD)
add_text(slide, 0.7, 1.4, 12, 0.4, 'RESPUESTA: SÍ, PERO PRIMERO HAY QUE LIMPIAR.', 22, True, RED, PP_ALIGN.CENTER)
add_text(slide, 0.7, 1.8, 12, 0.4, '6% de datos corruptos. Si se usa sin limpiar, el VP tomará decisiones erróneas.', 14, False, WHITE, PP_ALIGN.CENTER)

nums = [('15.000','Empleados','90% son Analistas\nPirámide muy plana'), ('1.417','En CAV','9,4% del total\nÁrea más grande'), ('6%','Datos Corruptos','81 Inconsistente\n69 negativos\n1.373 outliers')]
for i, (n, l, d) in enumerate(nums):
    x = 0.5 + i * 4.2
    add_box(slide, x, 2.6, 3.9, 2.3, CARD)
    add_text(slide, x, 2.7, 3.9, 0.6, n, 32, True, RED, PP_ALIGN.CENTER)
    add_text(slide, x, 3.2, 3.9, 0.4, l, 16, True, WHITE, PP_ALIGN.CENTER)
    add_text(slide, x, 3.7, 3.9, 1, d, 12, False, GRAY, PP_ALIGN.CENTER)

add_box(slide, 0.5, 5.3, 12.3, 1.8, CARD)
add_text(slide, 0.7, 5.4, 12, 0.3, 'RECOMENDACIÓN:', 14, True, RED)
add_text(slide, 0.7, 5.8, 12, 0.5, '1) Limpiar salarios  →  2) Unificar áreas  →  3) Resolver IDs duplicados  →  4) Auditar contra fuente real', 14, False, WHITE, PP_ALIGN.CENTER)
add_text(slide, 0.7, 6.3, 12, 0.5, '"Necesito decisiones ya. ¿Qué me recomiendas hacer con esta información?"\n→ LIMPIAR PRIMERO. Sin datos confiables no hay decisión estratégica válida.', 13, False, ORANGE, PP_ALIGN.CENTER)

# =============================================
# BLOQUE 4: ENTREVISTA PROFUNDA (NOTAS)
# =============================================
section_divider('BLOQUE 4', 'Entrevista Profunda — 30 min — Respuestas orales (no presentar)')

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_box(slide, 0, 0, 13.333, 0.08, RED)
add_text(slide, 0.5, 0.2, 12, 0.5, 'BLOQUE 4 — RESPUESTAS DE ENTREVISTA', 24, True, RED)
add_text(slide, 0.5, 0.6, 12, 0.3, 'NOTAS: Estas respuestas son ORALES, no se presentan en slides', 12, False, ORANGE)

qa = [
    ('¿Cómo decides cuándo un dataset no sirve?',
     'Cuando el % de datos corruptos/faltantes supera 10-15% en campos críticos y no hay fuente para imputar. Aquí 6% es borderline.'),
    ('¿Qué harías si negocio insiste en datos malos?',
     'Documentar riesgos por escrito, mostrar ejemplos concretos, ofrecer versión limpia como alternativa.'),
]
y = 1.0
for q, a in qa:
    add_box(slide, 0.5, y, 12.3, 0.8, CARD)
    add_text(slide, 0.7, y+0.02, 11.5, 0.3, 'P: ' + q, 12, True, RED)
    add_text(slide, 0.7, y+0.35, 11.5, 0.4, 'R: ' + a, 11, False, WHITE)
    y += 0.85

# Page 2 of Block 4
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_box(slide, 0, 0, 13.333, 0.08, RED)
add_text(slide, 0.5, 0.2, 12, 0.5, 'BLOQUE 4 — RESPUESTAS (cont.)', 24, True, RED)

qa2 = [
    ('¿Cómo modelarías este problema? (Técnica)',
     'Pipeline ETL: Extract (leer crudo) → Transform (limpiar, imputar, validar) → Load (carga limpia). Python/Pandas con reglas automatizadas.'),
    ('¿Qué técnicas para categóricos esparsos? (Técnica)',
     'One-Hot si baja cardinalidad, Target Encoding si alta, o agrupar raras en "Otros". Aquí: unificar nombres de áreas.'),
    ('¿Cómo impacta en decisiones organizacionales? (Negocio)',
     'Salario "Inconsistente" distorsiona presupuestos. ID duplicado sobredimensiona áreas. Datos faltantes impide análisis completo.'),
    ('¿Qué KPIs propondrías? (Negocio)',
     '1) % completitud, 2) Tasa duplicados, 3) Rango salarial válido, 4) Cobertura seguimiento, 5) Tiempo desde última actualización.'),
    ('¿Qué es peor: no tener datos o datos malos? (Criterio)',
     'Datos malos. Sin datos sabes que no sabes. Con datos malos crees saber y tomas decisiones equivocadas.'),
]
y = 1.0
for q, a in qa2:
    add_box(slide, 0.5, y, 12.3, 0.95, CARD)
    add_text(slide, 0.7, y+0.02, 11.5, 0.3, 'P: ' + q, 12, True, RED)
    add_text(slide, 0.7, y+0.35, 11.5, 0.55, 'R: ' + a, 11, False, WHITE)
    y += 1.0

# Save
out = '/Users/nicolassantos/Desktop/Proyectos/Assessment_Center_Analisis.pptx'
prs.save(out)
print(f'PPT: {out}')
print(f'Slides: {len(prs.slides)}')
